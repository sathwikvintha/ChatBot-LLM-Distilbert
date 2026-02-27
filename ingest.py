# Import standard libraries
from http import client
import os
import json
import logging
import time
from pathlib import Path
from typing import List, Dict
import re
import pdfplumber
import docx
from sentence_transformers import SentenceTransformer
import numpy as np
import pandas as pd
from pptx import Presentation
from qdrant_client import QdrantClient
from qdrant_client.models import Distance, VectorParams, PointStruct
import hashlib
import uuid
from qdrant_client.models import Filter, FieldCondition, MatchValue

from config import QDRANT_PATH, COLLECTION_NAME

# -----------------------------
# Setup logging
# -----------------------------
logging.basicConfig(
    level=logging.INFO,
    format="%(asctime)s [%(levelname)s] %(message)s",
    handlers=[
        logging.FileHandler("ingest.log"),
        logging.StreamHandler()
    ]
)

# -----------------------------
# Define directories
# -----------------------------
SOURCE_DIR = Path("data/source")
PROCESSED_DIR = Path("data/processed")
CHUNKS_DIR = Path("data/chunks")
EMBEDDINGS_DIR = Path("data/embeddings")

for d in [PROCESSED_DIR, CHUNKS_DIR, EMBEDDINGS_DIR]:
    d.mkdir(parents=True, exist_ok=True)

# -----------------------------
# Load embedding model
# -----------------------------
MODEL_NAME = "sentence-transformers/all-MiniLM-L6-v2"
embedder = SentenceTransformer(MODEL_NAME)

# -----------------------------
# Trash cleaner
# -----------------------------
def clean_trash(text: str) -> str:
    """
    Remove junk lines like '|' bars, repeated dashes, and UI noise.
    Keeps meaningful content, collapses symbol spam.
    """
    lines = text.splitlines()
    cleaned_lines = []
    for ln in lines:
        s = ln.strip()
        if not s:
            continue
        if re.fullmatch(r"[| \-]+", s):
            continue
        symbols = sum(1 for c in s if not c.isalnum())
        if symbols / max(1, len(s)) > 0.7:
            continue
        s = re.sub(r"(\|[ \|]*)+", "|", s)
        s = re.sub(r"(-{2,})", "-", s)
        s = re.sub(r"\s+", " ", s)
        cleaned_lines.append(s)
    return "\n".join(cleaned_lines)

# -----------------------------
# Text normalization
# -----------------------------
def normalize_text(text: str) -> str:
    text = re.sub(r"\r\n", "\n", text)
    text = re.sub(r"[ \t]+", " ", text)
    text = re.sub(r"(?<!\n)([-*•]\s)", r"\n\1", text)
    text = re.sub(r"(?<!\n)(\d+\.\s)", r"\n\1", text)
    return text.strip()

# -----------------------------
# Extraction functions (PDF, DOCX, TXT only)
# -----------------------------
def extract_pdf_text(path: Path) -> str:
    pages = []
    with pdfplumber.open(path) as pdf:
        for page in pdf.pages:
            text = page.extract_text() or ""
            text = clean_trash(text)
            pages.append(text)
    return "\n".join(pages)

def extract_docx_text(path: Path) -> str:
    doc = docx.Document(path)
    text = "\n".join(p.text for p in doc.paragraphs if p.text.strip())
    return clean_trash(text)

def extract_excel_text(path: Path) -> str:
    text_blocks = []
    excel_data = pd.read_excel(path, sheet_name=None)

    for sheet_name, df in excel_data.items():
        text_blocks.append(f"\nSheet: {sheet_name}")
        df = df.fillna("")  # remove NaNs
        text_blocks.append(df.astype(str).to_string(index=False))

    return clean_trash("\n".join(text_blocks))

def extract_ppt_text(path: Path) -> str:
    prs = Presentation(path)
    slides_text = []

    for i, slide in enumerate(prs.slides):
        slides_text.append(f"\nSlide {i + 1}")
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                slides_text.append(shape.text.strip())

    return clean_trash("\n".join(slides_text))

def extract_text(path: Path) -> str:
    suffix = path.suffix.lower()

    if suffix == ".pdf":
        return extract_pdf_text(path)

    elif suffix == ".docx":
        return extract_docx_text(path)

    elif suffix in [".txt", ".md"]:
        return clean_trash(path.read_text(encoding="utf-8", errors="ignore"))

    elif suffix in [".xlsx", ".xls"]:
        return extract_excel_text(path)

    elif suffix == ".pptx":
        return extract_ppt_text(path)

    else:
        raise ValueError(f"Unsupported file type: {path}")

def save_processed(path: Path, text: str):
    out = PROCESSED_DIR / f"{path.stem}.txt"
    out.write_text(text, encoding="utf-8")

# -----------------------------
# Chunking
# -----------------------------
def _split_blocks(text: str) -> List[str]:
    lines = [l.strip() for l in text.split("\n")]
    blocks, current = [], []
    def flush():
        if current:
            blocks.append("\n".join(current).strip())
            current.clear()
    for ln in lines:
        if not ln:
            flush()
            continue
        if re.match(r"^(\d+\.\s|[-*•]\s)", ln) or ln.startswith("mvn ") or ln.startswith("-D"):
            flush()
            blocks.append(ln)
        else:
            current.append(ln)
    flush()
    return [b for b in blocks if b and len(b) > 1]

def chunk_text(text: str, max_chars=450) -> List[str]:
    blocks = _split_blocks(text)
    chunks, buf = [], ""
    for b in blocks:
        if len(b) > max_chars:
            sentences = re.split(r'(?<=[.!?])\s+', b)
            cur = ""
            for s in sentences:
                if len(cur) + len(s) <= max_chars:
                    cur += (" " if cur else "") + s
                else:
                    if cur:
                        chunks.append(cur.strip())
                    cur = s
            if cur.strip():
                chunks.append(cur.strip())
            continue
        if len(buf) + len(b) <= max_chars:
            buf += ("\n" if buf else "") + b
        else:
            if buf.strip():
                chunks.append(buf.strip())
            buf = b
    if buf.strip():
        chunks.append(buf.strip())
    uniq, seen = [], set()
    for c in chunks:
        c_norm = re.sub(r"\s+", " ", c).strip()
        if len(c_norm) < 25:
            continue
        if c_norm.lower() in {"not mentioned in the document", "not relevant to the attached documents"}:
            continue
        if c_norm not in seen:
            uniq.append(c_norm)
            seen.add(c_norm)
    return uniq

def build_chunks(path: Path, text: str) -> List[Dict]:
    raw_chunks = chunk_text(text)
    return [
        {
            "chunk_id": f"{path.stem}_{i}",
            "source_path": str(path),
            "file_name": path.name,
            "doc_type": path.suffix.lower(),
            "file_size_kb": round(path.stat().st_size / 1024, 2),
            "indexed_at": time.time(),
            "text": chunk,
            "has_ocr": False
        }
        for i, chunk in enumerate(raw_chunks)
    ]

def save_chunks(path: Path, chunks: List[Dict]):
    out = CHUNKS_DIR / f"{path.stem}.jsonl"
    with out.open("w", encoding="utf-8") as f:
        for c in chunks:
            f.write(json.dumps(c, ensure_ascii=False) + "\n")

# -----------------------------
# Embeddings + index
# -----------------------------
def embed_chunks(chunks: List[Dict]) -> np.ndarray:
    texts = [c["text"] for c in chunks]
    return embedder.encode(texts, convert_to_numpy=True, normalize_embeddings=True)

def save_embeddings(path: Path, embeddings: np.ndarray):
    np.save(EMBEDDINGS_DIR / f"{path.stem}.npy", embeddings)
HASH_STORE = Path("data/file_hashes.json")

def compute_file_hash(path: Path) -> str:
    hash_md5 = hashlib.md5()
    with open(path, "rb") as f:
        for chunk in iter(lambda: f.read(8192), b""):
            hash_md5.update(chunk)
    return hash_md5.hexdigest()

def load_hash_store():
    if HASH_STORE.exists():
        return json.loads(HASH_STORE.read_text())
    return {}

def save_hash_store(hash_data):
    HASH_STORE.write_text(json.dumps(hash_data, indent=2))
    
# -----------------------------
# Main ingestion pipeline
# -----------------------------
def ingest():
    client = QdrantClient(path=QDRANT_PATH)
    collection_name = COLLECTION_NAME  

    if collection_name not in [c.name for c in client.get_collections().collections]:
        dim = embedder.get_sentence_embedding_dimension()
        client.create_collection(
            collection_name=collection_name,
            vectors_config=VectorParams(
                size=dim,
                distance=Distance.COSINE
            ),
        )
    hash_store = load_hash_store()  
    BATCH_SIZE = 500
    points_batch = []
    SUPPORTED_TYPES = {".pdf", ".docx", ".txt", ".md", ".xlsx", ".xls", ".pptx"}
    
    for file in SOURCE_DIR.iterdir():  
        if not file.is_file():
            continue
        if file.suffix.lower() not in SUPPORTED_TYPES:
            logging.warning(f"Skipping unsupported file: {file.name}")
            continue
        file_hash = compute_file_hash(file)
        if file.name in hash_store and hash_store[file.name] == file_hash:
            logging.info(f"Skipping already indexed file: {file.name}")
            continue
        if file.name in hash_store and hash_store[file.name] != file_hash:
            logging.info(f"File changed. Deleting old vectors for {file.name}")

            client.delete(
                collection_name=collection_name,
                points_selector=Filter(
                must=[
                    FieldCondition(
                        key="file_name",
                        match=MatchValue(value=file.name)
                    )
                ]
            )
        )
        t0 = time.time()

        try:
            logging.info(f"Processing {file.name}...")
            raw_text = extract_text(file)
            clean_text = normalize_text(raw_text)
            clean_text = clean_trash(clean_text)
            save_processed(file, clean_text)
            chunks = build_chunks(file, clean_text)
            save_chunks(file, chunks)
            embeddings = embed_chunks(chunks)

            for chunk, embedding in zip(chunks, embeddings):
                points_batch.append(
                    PointStruct(
                        id=str(uuid.uuid4()),
                        vector=embedding.tolist(),
                        payload=chunk
                    )
                )
                # Upload batch when size reached
                if len(points_batch) >= BATCH_SIZE:
                    client.upsert(
                        collection_name=collection_name,
                        points=points_batch
                    )
                    logging.info(f"Uploaded {len(points_batch)} vectors")
                    points_batch = []
            logging.info(
                f"{file.name}: {len(chunks)} chunks processed "
                f"in {time.time() - t0:.2f}s"
            )
            hash_store[file.name] = file_hash
            save_hash_store(hash_store)
        except Exception as e:
            logging.error(f"Failed to process {file.name}: {e}")
            continue

    # Upload remaining points
    if points_batch:
        client.upsert(
            collection_name=collection_name,
            points=points_batch
        )
        logging.info(f"Uploaded final {len(points_batch)} vectors")
    logging.info("Ingestion complete. Qdrant index ready.")

if __name__ == "__main__":
    ingest()
